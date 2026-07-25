"""
Executive Presentation Generator
Generates 5-7 slide executive presentations from project health data
"""

import json
import os
import io
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


CATEGORY_ORDER = ["Schedule", "Blockages", "Stakeholder Attitude", "Budget Burn"]
CATEGORY_WEIGHTS = {
    "Schedule": 0.25,
    "Budget Burn": 0.25,
    "Blockages": 0.25,
    "Stakeholder Attitude": 0.25,
}
WEIGHTED_STATUS_PCT = {"Green": 100, "Amber": 50, "Red": 0}
RAG_STATUS_VALUES = {"Green": 3, "Amber": 2, "Red": 1}


def _clean_reasoning(reasoning, max_len=180):
    """Strip markdown/artifacts from the reasoning string and truncate cleanly."""
    if not reasoning:
        return "No reasoning available"
    cleaned = reasoning
    for prefix in ["**Overall Status: Green**", "**Overall Status: Amber**", "**Overall Status: Red**", "Overall Status: Green", "Overall Status: Amber", "Overall Status: Red"]:
        if cleaned.startswith(prefix):
            cleaned = cleaned[len(prefix):]
    cleaned = cleaned.replace("**", "")
    cleaned = cleaned.replace("ò", "•").replace("• ", "•").replace("••", "•")
    cleaned = cleaned.replace("\n", " ").replace("\r", " ").strip()
    while "  " in cleaned:
        cleaned = cleaned.replace("  ", " ")
    if len(cleaned) > max_len:
        cut = cleaned[:max_len]
        last_space = cut.rfind(" ")
        if last_space > max_len * 0.7:
            cut = cut[:last_space]
        cleaned = cut + "…"
    return cleaned.strip("• :,-")


def _normalize_status(status):
    if not status or not isinstance(status, str):
        return None
    s = status.strip().title()
    return s if s in RAG_STATUS_VALUES else None


def _compute_portfolio_metrics(projects):
    """Compute the same health metrics as the dashboard (weighted score + rollups)."""
    total_projects = len(projects or [])
    green_count = sum(1 for p in projects if _normalize_status(p.get('overall_status')) == 'Green')
    amber_count = sum(1 for p in projects if _normalize_status(p.get('overall_status')) == 'Amber')
    red_count = sum(1 for p in projects if _normalize_status(p.get('overall_status')) == 'Red')

    green_categories = 0
    total_categories = 0
    earned_points = 0
    max_points = 0
    project_weighted_scores = []

    for project in projects:
        cat_scores = project.get('category_scores', {}) or {}
        weighted_total = 0.0
        for category in CATEGORY_ORDER:
            status = _normalize_status(cat_scores.get(category))
            if not status:
                continue
            total_categories += 1
            max_points += RAG_STATUS_VALUES['Green']
            earned_points += RAG_STATUS_VALUES[status]
            if status == 'Green':
                green_categories += 1
            weighted_total += CATEGORY_WEIGHTS.get(category, 0.0) * WEIGHTED_STATUS_PCT[status]
        project_weighted_scores.append(weighted_total)

    project_rollup = int((green_count / total_projects) * 100) if total_projects else 0
    category_rollup = int((green_categories / total_categories) * 100) if total_categories else 0
    health_index = int((earned_points / max_points) * 100) if max_points else 0
    weighted = int(sum(project_weighted_scores) / len(project_weighted_scores)) if project_weighted_scores else 0

    return {
        'total_projects': total_projects,
        'green_count': green_count,
        'amber_count': amber_count,
        'red_count': red_count,
        'project_rollup_pct': project_rollup,      # fully-green projects
        'category_rollup_pct': category_rollup,    # green categories
        'health_index_pct': health_index,          # earned RAG points
        'weighted_pct': weighted,                  # dashboard KPI (primary)
        'green_categories': green_categories,
        'total_categories': total_categories,
    }


class PresentationGenerator:
    """Generate executive presentations from project health data"""

    def __init__(self, output_dir="presentations"):
        self.output_dir = output_dir
        self._dir_created = False
        import tempfile
        candidates = [output_dir, os.path.join(tempfile.gettempdir(), "phr_presentations")]
        for candidate in candidates:
            try:
                if not os.path.exists(candidate):
                    os.makedirs(candidate, exist_ok=True)
                self.output_dir = candidate
                self._dir_created = True
                break
            except Exception:
                continue

    def load_project_data(self, uploaded_files=None):
        """Load project health data (prefer uploaded_files, fallback to projects/ dir, weekly_outputs)."""
        # (A) uploaded files take precedence
        if uploaded_files:
            from project_health_agent import ProjectHealthAgent, load_project_from_excel
            import tempfile

            agent = ProjectHealthAgent()
            projects_data = {}

            for uploaded_file in uploaded_files:
                tmp_path = None
                try:
                    try:
                        uploaded_file.seek(0)
                    except Exception:
                        pass
                    suffix = os.path.splitext(getattr(uploaded_file, 'name', '') or '')[1] or '.xlsx'
                    with tempfile.NamedTemporaryFile(
                        delete=False,
                        suffix=suffix,
                        dir=tempfile.gettempdir(),
                        prefix='phr_gen_'
                    ) as tmp:
                        tmp.write(uploaded_file.getbuffer())
                        tmp.flush()
                        try:
                            os.fsync(tmp.fileno())
                        except Exception:
                            pass
                        tmp_path = tmp.name

                    try:
                        df = load_project_from_excel(tmp_path)
                        if df is not None and not df.empty:
                            project_name = os.path.splitext(uploaded_file.name)[0]
                            result = agent.evaluate_project(df)
                            projects_data[project_name] = {
                                'project_name': project_name,
                                'overall_status': result.status,
                                'reasoning': result.reasoning,
                                'category_scores': result.category_scores,
                                'category_details': result.details,
                                'timestamp': datetime.now().isoformat()
                            }
                    except Exception as e:
                        print(f"Error evaluating {uploaded_file.name}: {e}")
                        continue
                except Exception as e:
                    print(f"Error loading {uploaded_file.name}: {e}")
                finally:
                    if tmp_path and os.path.exists(tmp_path):
                        try:
                            os.unlink(tmp_path)
                        except Exception:
                            pass

            if projects_data:
                return list(projects_data.values())

        # (B) Fallback 1 — projects/ directory on disk
        projects_dir = "projects"
        if os.path.exists(projects_dir):
            from project_health_agent import ProjectHealthAgent, load_project_from_excel
            agent = ProjectHealthAgent()
            results = []
            try:
                filenames = sorted(os.listdir(projects_dir))
            except Exception:
                filenames = []
            for fn in filenames:
                if not fn.lower().endswith(('.xlsx', '.xls', '.csv')):
                    continue
                fp = os.path.join(projects_dir, fn)
                try:
                    df = load_project_from_excel(fp)
                    if df is not None and not df.empty:
                        project_name = os.path.splitext(fn)[0]
                        res = agent.evaluate_project(df)
                        results.append({
                            'project_name': project_name,
                            'overall_status': res.status,
                            'reasoning': res.reasoning,
                            'category_scores': res.category_scores,
                            'category_details': res.details,
                            'timestamp': datetime.now().isoformat(),
                        })
                except Exception as e:
                    print(f"Error loading {fp}: {e}")
            if results:
                return results

        # (C) Fallback 2 — weekly_outputs JSON reports
        weekly_dir = "weekly_outputs"
        if os.path.exists(weekly_dir):
            try:
                all_entries = os.listdir(weekly_dir)
            except Exception:
                all_entries = []
            reports = sorted(
                [os.path.join(weekly_dir, f) for f in all_entries if f.endswith('.json') and not f.startswith('summary_')],
                key=lambda p: os.path.getmtime(p),
                reverse=True,
            )
            results = []
            seen_projects = set()
            for rpath in reports:
                try:
                    with open(rpath, 'r', encoding='utf-8') as f:
                        data = json.load(f)
                except Exception:
                    continue
                name = (data.get('project_name') or os.path.splitext(os.path.basename(rpath))[0]).strip()
                if name in seen_projects:
                    continue
                seen_projects.add(name)
                results.append({
                    'project_name': name,
                    'overall_status': data.get('overall_status') or data.get('status') or 'Unknown',
                    'reasoning': data.get('reasoning') or data.get('summary') or '',
                    'category_scores': data.get('category_scores') or data.get('scores') or {},
                    'category_details': data.get('category_details') or data.get('details') or {},
                    'timestamp': data.get('timestamp') or datetime.now().isoformat(),
                })
            if results:
                return results

        return None

    def create_presentation(self, projects):
        """Generate executive presentation"""
        if not projects:
            print("No project data available")
            return None

        prs = Presentation()

        metrics = _compute_portfolio_metrics(projects)
        total_projects = metrics['total_projects']
        green_count = metrics['green_count']
        amber_count = metrics['amber_count']
        red_count = metrics['red_count']
        health_score = metrics['weighted_pct']

        # Slide 1: Title Slide
        self._add_title_slide(
            prs,
            "Project Health Executive Summary",
            f"Portfolio Health Score: {health_score}% | {datetime.now().strftime('%B %d, %Y')}"
        )

        # Slide 2: Portfolio Overview
        self._add_portfolio_overview(prs, metrics)

        # Slide 3: Key Risks & Highlights
        self._add_key_risks(prs, projects)

        # Slide 4: Project Details
        self._add_project_details(prs, projects)

        # Slide 5: Trends & Insights
        self._add_trends_insights(prs, projects, metrics)

        # Slide 6: Recommendations
        self._add_recommendations(prs, projects)

        # Slide 7: Next Steps
        self._add_next_steps(prs)

        return prs
    
    def _add_title_slide(self, prs, title, subtitle):
        """Add title slide"""
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        
        title_shape.text = title
        subtitle_shape.text = subtitle
        
        # Style title
        title_frame = title_shape.text_frame
        title_frame.paragraphs[0].font.size = Pt(44)
        title_frame.paragraphs[0].font.color.rgb = RGBColor(10, 25, 47)  # Navy
        
        # Style subtitle
        subtitle_frame = subtitle_shape.text_frame
        subtitle_frame.paragraphs[0].font.size = Pt(24)
        subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(0, 180, 216)  # Teal
    
    def _add_portfolio_overview(self, prs, metrics):
        """Add portfolio overview slide"""
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)

        title_shape = slide.shapes.title
        title_shape.text = "Portfolio Overview"

        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True

        total = metrics['total_projects']
        green = metrics['green_count']
        amber = metrics['amber_count']
        red = metrics['red_count']
        weighted = metrics['weighted_pct']
        health_index = metrics['health_index_pct']
        cat_rollup = metrics['category_rollup_pct']
        fully_green_pct = metrics['project_rollup_pct']

        p = text_frame.add_paragraph()
        p.text = f"Portfolio Health Score (Weighted KPI): {weighted}%"
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(10, 25, 47)
        p.space_after = Pt(6)

        p = text_frame.add_paragraph()
        p.text = f"Health Index: {health_index}%  ·  Green Parameters: {cat_rollup}%  ·  Fully-Green Projects: {fully_green_pct}%"
        p.font.size = Pt(14)
        p.font.italic = True
        p.font.color.rgb = RGBColor(80, 80, 80)
        p.space_after = Pt(14)

        p = text_frame.add_paragraph()
        p.text = f"Total Projects: {total}"
        p.font.size = Pt(20)
        p.space_after = Pt(8)

        p = text_frame.add_paragraph()
        p.text = f"• Green (On Track): {green}"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(40, 167, 69)
        p.space_after = Pt(4)

        p = text_frame.add_paragraph()
        p.text = f"• Amber (Monitor): {amber}"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(255, 193, 7)
        p.space_after = Pt(4)

        p = text_frame.add_paragraph()
        p.text = f"• Red (Critical): {red}"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(220, 53, 69)
        p.space_after = Pt(12)

        p = text_frame.add_paragraph()
        if weighted >= 75:
            verdict = "HEALTHY - The weighted dashboard KPI indicates strong portfolio performance."
        elif weighted >= 50:
            verdict = "MODERATE RISK - Warning indicators are present; proactive management required."
        else:
            verdict = "HIGH RISK - Critical issues detected; immediate intervention required."
        p.text = f"Portfolio Status: {verdict}"
        p.font.size = Pt(18)
        p.font.italic = True
        p.font.color.rgb = RGBColor(100, 100, 100)

    def _add_key_risks(self, prs, projects):
        """Add key risks and highlights slide"""
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)

        title_shape = slide.shapes.title
        title_shape.text = "Key Risks & Highlights"

        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True

        red_projects = [p for p in projects if _normalize_status(p.get('overall_status')) == 'Red']
        amber_projects = [p for p in projects if _normalize_status(p.get('overall_status')) == 'Amber']
        green_projects = [p for p in projects if _normalize_status(p.get('overall_status')) == 'Green']

        if red_projects:
            p = text_frame.add_paragraph()
            p.text = "CRITICAL RISKS (Immediate Attention Required):"
            p.font.size = Pt(22)
            p.font.bold = True
            p.font.color.rgb = RGBColor(220, 53, 69)
            p.space_after = Pt(10)

            for project in red_projects:
                p = text_frame.add_paragraph()
                p.text = f"• {project.get('project_name', 'Unknown')}: {_clean_reasoning(project.get('reasoning'))}"
                p.font.size = Pt(15)
                p.space_after = Pt(8)
                p.level = 1

        if amber_projects:
            p = text_frame.add_paragraph()
            p.text = "MONITORING REQUIRED:"
            p.font.size = Pt(22)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 193, 7)
            p.space_after = Pt(10)

            for project in amber_projects:
                p = text_frame.add_paragraph()
                p.text = f"• {project.get('project_name', 'Unknown')}: {_clean_reasoning(project.get('reasoning'))}"
                p.font.size = Pt(15)
                p.space_after = Pt(8)
                p.level = 1

        if green_projects:
            p = text_frame.add_paragraph()
            p.text = "HIGHLIGHTS (On Track):"
            p.font.size = Pt(22)
            p.font.bold = True
            p.font.color.rgb = RGBColor(40, 167, 69)
            p.space_after = Pt(10)
            for project in green_projects:
                p = text_frame.add_paragraph()
                p.text = f"• {project.get('project_name', 'Unknown')}: {_clean_reasoning(project.get('reasoning'), max_len=120)}"
                p.font.size = Pt(15)
                p.space_after = Pt(6)
                p.level = 1

        if not red_projects and not amber_projects and not green_projects:
            p = text_frame.add_paragraph()
            p.text = "No project data available for risk analysis."
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(120, 120, 120)

    def _add_project_details(self, prs, projects):
        """Add project details slide"""
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)

        title_shape = slide.shapes.title
        title_shape.text = "Project Status Details"

        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True

        for project in projects:
            overall = _normalize_status(project.get('overall_status')) or project.get('overall_status', 'Unknown')
            p = text_frame.add_paragraph()
            p.text = f"{project.get('project_name', 'Unknown')} — {overall}"
            p.font.size = Pt(18)
            p.font.bold = True
            if overall == 'Green':
                p.font.color.rgb = RGBColor(40, 167, 69)
            elif overall == 'Amber':
                p.font.color.rgb = RGBColor(200, 150, 0)
            elif overall == 'Red':
                p.font.color.rgb = RGBColor(220, 53, 69)
            p.space_after = Pt(6)

            p = text_frame.add_paragraph()
            p.text = _clean_reasoning(project.get('reasoning'), max_len=220)
            p.font.size = Pt(12)
            p.font.italic = True
            p.font.color.rgb = RGBColor(90, 90, 90)
            p.space_after = Pt(4)
            p.level = 1

            category_scores = project.get('category_scores', {}) or {}
            for category in CATEGORY_ORDER:
                score = _normalize_status(category_scores.get(category))
                if not score:
                    continue
                p = text_frame.add_paragraph()
                p.text = f"  {category}: {score}"
                p.font.size = Pt(13)
                p.space_after = Pt(2)
                p.level = 2

            p.space_after = Pt(10)

    def _add_trends_insights(self, prs, projects, metrics=None):
        """Add trends and insights slide"""
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)

        title_shape = slide.shapes.title
        title_shape.text = "Trends & Insights"

        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True

        schedule_issues = sum(1 for p in projects if _normalize_status((p.get('category_scores') or {}).get('Schedule')) in ['Red', 'Amber'])
        budget_issues = sum(1 for p in projects if _normalize_status((p.get('category_scores') or {}).get('Budget Burn')) in ['Red', 'Amber'])
        blockage_issues = sum(1 for p in projects if _normalize_status((p.get('category_scores') or {}).get('Blockages')) in ['Red', 'Amber'])
        stakeholder_issues = sum(1 for p in projects if _normalize_status((p.get('category_scores') or {}).get('Stakeholder Attitude')) in ['Red', 'Amber'])

        if metrics is not None:
            p = text_frame.add_paragraph()
            p.text = f"PORTFOLIO AT A GLANCE: {metrics['total_projects']} projects · Weighted KPI {metrics['weighted_pct']}% · Category Rollup {metrics['category_rollup_pct']}%"
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(60, 60, 60)
            p.space_after = Pt(10)

        p = text_frame.add_paragraph()
        p.text = "KEY OBSERVATIONS:"
        p.font.size = Pt(22)
        p.font.bold = True
        p.space_after = Pt(10)

        if schedule_issues > 0:
            p = text_frame.add_paragraph()
            p.text = f"• Schedule: {schedule_issues} project(s) behind baseline. Timeline execution is the #1 driver of current RAG status."
            p.font.size = Pt(15)
            p.space_after = Pt(6)
        else:
            p = text_frame.add_paragraph()
            p.text = "• Schedule: All projects tracking to baseline (Green across portfolio)."
            p.font.size = Pt(15)
            p.space_after = Pt(6)

        if budget_issues > 0:
            p = text_frame.add_paragraph()
            p.text = f"• Budget: {budget_issues} project(s) with burn variance — financial oversight recommended."
            p.font.size = Pt(15)
            p.space_after = Pt(6)
        else:
            p = text_frame.add_paragraph()
            p.text = "• Budget: Spend aligned with completion across the portfolio."
            p.font.size = Pt(15)
            p.space_after = Pt(6)

        if blockage_issues > 0:
            p = text_frame.add_paragraph()
            p.text = f"• Blockages: {blockage_issues} project(s) with active blockers — escalation paths should be verified."
            p.font.size = Pt(15)
            p.space_after = Pt(6)
        else:
            p = text_frame.add_paragraph()
            p.text = "• Blockages: No material blockers reported."
            p.font.size = Pt(15)
            p.space_after = Pt(6)

        if stakeholder_issues > 0:
            p = text_frame.add_paragraph()
            p.text = f"• Stakeholders: {stakeholder_issues} project(s) with attitude flags — communication cadence should increase."
            p.font.size = Pt(15)
            p.space_after = Pt(6)
        else:
            p = text_frame.add_paragraph()
            p.text = "• Stakeholders: Sentiment positive or neutral (no news = good news)."
            p.font.size = Pt(15)
            p.space_after = Pt(6)

        p = text_frame.add_paragraph()
        p.text = "INSIGHT: Portfolio health is driven primarily by schedule adherence. Focus recovery effort first on tasks with late critical-path milestones and validate that Duration and % Complete fields are being updated weekly."
        p.font.size = Pt(14)
        p.font.italic = True
        p.font.color.rgb = RGBColor(80, 80, 80)
        p.space_after = Pt(10)
    
    def _add_recommendations(self, prs, projects):
        """Add recommendations slide"""
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)

        title_shape = slide.shapes.title
        title_shape.text = "Strategic Recommendations"

        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True

        red_projects = [p for p in projects if _normalize_status(p.get('overall_status')) == 'Red']
        amber_projects = [p for p in projects if _normalize_status(p.get('overall_status')) == 'Amber']

        p = text_frame.add_paragraph()
        p.text = "IMMEDIATE ACTIONS (This Week):"
        p.font.size = Pt(22)
        p.font.bold = True
        p.space_after = Pt(10)

        if red_projects:
            for project in red_projects:
                p = text_frame.add_paragraph()
                p.text = f"• {project.get('project_name', 'Unknown')}: Activate recovery plan; re-plan critical-path baseline; daily stand-up until caught up."
                p.font.size = Pt(15)
                p.space_after = Pt(6)

        if amber_projects:
            for project in amber_projects:
                p = text_frame.add_paragraph()
                p.text = f"• {project.get('project_name', 'Unknown')}: Add schedule buffer to slipping milestones, confirm stakeholder commitment."
                p.font.size = Pt(15)
                p.space_after = Pt(6)

        if not red_projects and not amber_projects:
            p = text_frame.add_paragraph()
            p.text = "• Continue current practices — portfolio is healthy; no immediate interventions required."
            p.font.size = Pt(15)
            p.space_after = Pt(6)

        p = text_frame.add_paragraph()
        p.text = "ONGOING MONITORING:"
        p.font.size = Pt(22)
        p.font.bold = True
        p.space_after = Pt(10)

        p = text_frame.add_paragraph()
        p.text = "• Weekly health reviews for all projects with updated Duration, % Complete, and Start/End dates"
        p.font.size = Pt(15)
        p.space_after = Pt(6)

        p = text_frame.add_paragraph()
        p.text = "• Proactive stakeholder communication for Amber/Red projects (bi-weekly minimum)"
        p.font.size = Pt(15)
        p.space_after = Pt(6)

        p = text_frame.add_paragraph()
        p.text = "• Resource allocation adjustments based on schedule priority and baseline gaps"
        p.font.size = Pt(15)
        p.space_after = Pt(6)
    
    def _add_next_steps(self, prs):
        """Add next steps slide"""
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        title_shape.text = "Next Steps"
        
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        
        p = text_frame.add_paragraph()
        p.text = "IMMEDIATE (This Week):"
        p.font.size = Pt(22)
        p.font.bold = True
        p.space_after = Pt(10)
        
        p = text_frame.add_paragraph()
        p.text = "• Review critical project recovery plans"
        p.font.size = Pt(16)
        p.space_after = Pt(6)
        
        p = text_frame.add_paragraph()
        p.text = "• Schedule stakeholder update meetings"
        p.font.size = Pt(16)
        p.space_after = Pt(12)
        
        p = text_frame.add_paragraph()
        p.text = "SHORT-TERM (Next 2 Weeks):"
        p.font.size = Pt(22)
        p.font.bold = True
        p.space_after = Pt(10)
        
        p = text_frame.add_paragraph()
        p.text = "• Implement recommended interventions"
        p.font.size = Pt(16)
        p.space_after = Pt(6)
        
        p = text_frame.add_paragraph()
        p.text = "• Monitor progress on action items"
        p.font.size = Pt(16)
        p.space_after = Pt(6)
        
        p = text_frame.add_paragraph()
        p.text = "• Update portfolio health metrics"
        p.font.size = Pt(16)
        p.space_after = Pt(12)
        
        p = text_frame.add_paragraph()
        p.text = "Follow-up presentation scheduled for: " + datetime.now().strftime('%B %d, %Y')
        p.font.size = Pt(14)
        p.font.italic = True
        p.font.color.rgb = RGBColor(100, 100, 100)
    
    def presentation_to_bytes(self, prs):
        """Serialize a Presentation to in-memory bytes (no disk I/O).

        Returns a tuple of (bytes, suggested_filename) suitable for
        ``st.download_button(data=..., file_name=...)``.
        """
        if prs is None:
            return None, None
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        timestamp = datetime.now().strftime('%Y-%m-%d_%H-%M-%S')
        filename = f"executive_summary_{timestamp}.pptx"
        return buf.getvalue(), filename

    def save_presentation(self, prs, filename=None):
        """Save presentation to file (best-effort; never raises — in-memory path is primary).

        Returns the filepath on success, or ``None`` if the filesystem is not
        writable (the caller should fall back to ``presentation_to_bytes``).
        """
        if prs is None:
            return None
        if not filename:
            timestamp = datetime.now().strftime('%Y-%m-%d_%H-%M-%S')
            filename = f"executive_summary_{timestamp}.pptx"

        if not self._dir_created:
            try:
                os.makedirs(self.output_dir, exist_ok=True)
                self._dir_created = True
            except Exception:
                return None

        filepath = os.path.join(self.output_dir, filename)
        try:
            prs.save(filepath)
            print(f"Presentation saved to: {filepath}")
            return filepath
        except Exception:
            return None


def main():
    """Main execution function"""
    generator = PresentationGenerator()
    
    # Load project data
    projects = generator.load_project_data()
    
    if not projects:
        print("No project data found. Please run the agent first.")
        return
    
    # Generate presentation
    prs = generator.create_presentation(projects)
    
    if prs:
        # Save presentation (best-effort disk write)
        filepath = generator.save_presentation(prs)
        if filepath:
            print(f"Executive presentation generated successfully!")
            print(f"Total projects analyzed: {len(projects)}")
        else:
            # Fallback: confirm generation via in-memory serialization
            pptx_bytes, filename = generator.presentation_to_bytes(prs)
            if pptx_bytes is not None:
                print(f"Executive presentation generated successfully ({len(pptx_bytes)} bytes in memory, filename={filename}).")
                print(f"Total projects analyzed: {len(projects)}")
                print("Note: Filesystem is read-only — presentation served via in-memory download.")
            else:
                print("Failed to serialize presentation.")


if __name__ == "__main__":
    main()
